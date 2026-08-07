#!/usr/bin/env python3
"""Convert a PowerPoint deck to a Marp markdown twin.

Usage:
    python scripts/pptx_to_marp.py "<deck.pptx>" "<output lecture dir>" \
        --title "Structure" --week 2 --topic structure

Writes <outdir>/slides.md; embedded pictures go to <outdir>/img/.
Extraction is text-faithful, not layout-faithful (see design spec).
"""
import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

# Image formats Chromium (and therefore Marp's HTML/PDF/PPTX renderers) can
# display natively. Anything else (e.g. legacy Windows Metafile) is still
# extracted to disk, but flagged so the deck doesn't silently render a
# broken image.
WEB_SAFE_EXTS = {"png", "jpg", "jpeg", "gif", "webp", "svg", "bmp"}

# Fonts these decks use to mean "this line is source code", regardless of
# whether PowerPoint's own bullet/no-bullet paragraph property is set.
CODE_FONTS = {"courier new", "consolas", "courier", "menlo", "monaco",
              "lucida console"}

# Prefixes that mark a line as Java even when it carries none of the
# punctuation signals below (e.g. a lone "import java.util.List" with no
# trailing semicolon captured in its own paragraph).
CODE_KEYWORD_PREFIXES = (
    "public ", "private ", "protected ", "class ", "interface ", "enum ",
    "int ", "int[", "double ", "double[", "float ", "long ", "short ",
    "byte ", "boolean ", "char ", "char[", "String ", "String[", "void ",
    "return ", "return;", "static ", "final ", "import ", "package ",
    "new ", "if ", "if(", "else", "for ", "for(", "while ", "while(",
    "try ", "try{", "catch ", "catch(", "System.", "@Override",
)

# Non-breaking space / em space: several decks paste code with these in
# place of ASCII spaces (copy-pasted from rich text). Normalized only for
# the code-likeness *test*; the emitted fence line keeps the original text.
NBSP = "\xa0"
EMSP = " "


def _run_markdown(run):
    """A run's text, as a Markdown link if the run carries a hyperlink.

    Skipped when the visible text already IS the address (a bare URL is
    already a fine link target for Marp/GFM renderers).
    """
    text = run.text
    try:
        address = run.hyperlink.address
    except (AttributeError, KeyError):
        address = None
    if address and text and text != address:
        return f"[{text}]({address})"
    return text


def _para_text(para):
    """Join a paragraph's runs into one string, preserving hyperlinks."""
    return "".join(_run_markdown(run) for run in para.runs)


def _paragraph_default_font(para):
    """The paragraph's own default run typeface, if it sets one.

    Reads <a:pPr><a:defRPr><a:latin typeface="..."/>, which acts as a
    fallback for any run in the paragraph that doesn't override its font.
    """
    pPr = para._p.find(qn("a:pPr"))
    if pPr is None:
        return None
    defRPr = pPr.find(qn("a:defRPr"))
    if defRPr is None:
        return None
    latin = defRPr.find(qn("a:latin"))
    return latin.get("typeface") if latin is not None else None


def _is_code_font_paragraph(para):
    """True only if EVERY run carrying visible text resolves to a
    monospace/code typeface (own <a:rPr><a:latin>, falling back to the
    paragraph's <a:pPr><a:defRPr><a:latin> default).

    Deliberately whole-paragraph, not "any run": several decks style a
    single word as inline code inside an ordinary prose sentence (e.g.
    "...the deposit() and withdraw() methods encapsulate..."), and that
    must stay a bullet, not become a fence. A real code line has every
    run in the code font; an inline-styled word in a prose sentence does
    not.
    """
    default_font = _paragraph_default_font(para)
    runs_with_text = [r for r in para.runs if r.text.strip()]
    if not runs_with_text:
        return False
    for run in runs_with_text:
        rPr = run._r.find(qn("a:rPr"))
        latin = rPr.find(qn("a:latin")) if rPr is not None else None
        font = (latin.get("typeface") if latin is not None else None) or default_font
        if not font or font.strip().lower() not in CODE_FONTS:
            return False
    return True


def _looks_like_code(text):
    """Secondary confirmation: does the (already font-flagged) text also
    read like a line of Java, rather than prose/labels that merely
    borrowed the code font (e.g. a bare array-index diagram label)?
    """
    t = text.strip().replace(NBSP, " ").replace(EMSP, " ")
    if not t:
        return False
    if ";" in t or t.endswith("{") or t.startswith("}"):
        return True
    if t.startswith(("//", "/*", "*")):
        return True
    if " = " in t or " == " in t:
        return True
    if t.startswith(CODE_KEYWORD_PREFIXES):
        return True
    if re.search(r"\w\([^()]*\)", t):  # e.g. println(...), deposit() -- no
        return True                    # space before "(": excludes prose
    return False                       # like "...applies to (Cat, Tiger)"


def _is_code_paragraph(para):
    return _is_code_font_paragraph(para) and _looks_like_code(_para_text(para))


def text_frame_body_lines(text_frame):
    """Render one shape's paragraphs as Markdown body lines.

    Ordinary paragraphs become bullets (existing behaviour). Runs of
    consecutive whole-line-code paragraphs (see _is_code_paragraph) are
    grouped into a single fenced ```java block instead -- a shape with a
    5-line method body produces one fence, not 5 stray bullets.
    """
    lines = []
    code_buf = []

    def flush():
        if code_buf:
            lines.append("```java")
            lines.extend(code_buf)
            lines.append("```")
            code_buf.clear()

    for para in text_frame.paragraphs:
        text = _para_text(para)
        if not text.strip():
            continue
        if _is_code_paragraph(para):
            # Keep leading indentation (nested braces read as nested code);
            # only trailing whitespace is noise.
            code_buf.append(text.rstrip())
        else:
            flush()
            level = para.level or 0
            lines.append("  " * level + "- " + text.strip())
    flush()
    return lines


def sanitize_cell(text):
    """Collapse embedded line breaks into a single <br> and escape pipes.

    python-pptx joins a cell's paragraphs with "\\n" and represents a
    manual (soft) line break within a paragraph as "\\x0b" (vertical tab).
    Left raw, either character splits a `| ... |` table row across physical
    lines and corrupts the generated Markdown/Marp table.
    """
    text = re.sub(r"[ \t]*[\n\x0b]+[ \t\n\x0b]*", "<br>", text.strip())
    return text.replace("|", "\\|")


def table_markdown(table):
    rows = [[sanitize_cell(cell.text) for cell in row.cells]
            for row in table.rows]
    if not rows:
        return []
    out = ["| " + " | ".join(rows[0]) + " |",
           "|" + "---|" * len(rows[0])]
    out += ["| " + " | ".join(r) + " |" for r in rows[1:]]
    return out


# Namespace-qualified names for locating fill-embedded pictures: an
# <a:blip r:embed="rIdN"/> nested anywhere inside a NON-picture shape's XML
# (a picture placeholder left at its layout type, an autoshape/logo filled
# with an image, or a MEDIA shape's poster-frame blip) renders a real
# picture that MSO_SHAPE_TYPE.PICTURE never catches -- python-pptx
# classifies a shape as PLACEHOLDER or MEDIA before it ever considers
# whether the shape's own fill happens to be a picture.
BLIP_QN = qn("a:blip")
EMBED_QN = qn("r:embed")
CNVPR_QN = qn("p:cNvPr")


def _sniff_ext(blob):
    """Best-effort image format from magic bytes.

    Only reached when an image part's own partname carries no extension
    (not observed in this corpus, but the OPC spec doesn't forbid it).
    Falls back to 'png' -- the same default WEB_SAFE_EXTS already treats
    as safe, so an unrecognized format doesn't spuriously get flagged.
    """
    if blob.startswith(b"\x89PNG\r\n\x1a\n"):
        return "png"
    if blob.startswith(b"\xff\xd8\xff"):
        return "jpg"
    if blob.startswith((b"GIF87a", b"GIF89a")):
        return "gif"
    if blob.startswith(b"BM"):
        return "bmp"
    if blob.startswith(b"RIFF") and blob[8:12] == b"WEBP":
        return "webp"
    return "png"


def _shape_alt_text(shape):
    """Alt text for a fill-embedded image: the shape's own descr (its
    <p:cNvPr descr="...">, the same source picture alt text is read from),
    falling back to the shape's name when the source deck set no
    description.
    """
    descr = ""
    cNvPr = shape._element.find(".//" + CNVPR_QN)
    if cNvPr is not None:
        descr = cNvPr.get("descr", "") or ""
    descr = " ".join(descr.split())  # collapse embedded newlines
    return descr or (shape.name or "")


def _extract_fill_images(shape, slide, idx, img_dir, n_img, rid_to_name,
                          images, image_alt, non_web_images,
                          unresolved_fills):
    """Extract picture(s) embedded as a FILL of `shape` (as opposed to a
    MSO_SHAPE_TYPE.PICTURE shape, which the caller already handles).

    Mutates images/image_alt/non_web_images/unresolved_fills in place and
    returns the updated n_img counter, so filenames keep continuing the
    same slideNN-M.<ext> sequence the PICTURE branch uses. A shape with no
    fill blip is a no-op. Dedupes by rId within the slide: a second shape
    whose blip carries an rId already written by an earlier shape on this
    slide reuses that file (no re-write, no re-warn) but still gets its
    own emitted image reference, since it's a distinct shape.
    """
    for blip in shape._element.findall(".//" + BLIP_QN):
        rid = blip.get(EMBED_QN)
        if not rid:
            continue  # e.g. a linked (r:link) blip, not an embedded one
        name = rid_to_name.get(rid)
        if name is None:
            try:
                part = slide.part.related_part(rid)
            except KeyError:
                unresolved_fills.append((shape.name, rid))
                continue
            ext = part.partname.ext or _sniff_ext(part.blob)
            n_img += 1
            name = f"slide{idx:02d}-{n_img}.{ext}"
            img_dir.mkdir(parents=True, exist_ok=True)
            (img_dir / name).write_bytes(part.blob)
            rid_to_name[rid] = name
            if ext.lower() not in WEB_SAFE_EXTS:
                non_web_images.append(name)
        images.append(name)
        image_alt[name] = _shape_alt_text(shape)
    return n_img


def extract_slide(slide, idx, img_dir):
    """Return (title, body_lines, image_names, notes, non_web_images,
    image_alt, media_count, unresolved_fills).

    non_web_images is the subset of image_names saved in a format that
    browsers -- and therefore Marp's Chromium-based renderers -- cannot
    display natively (e.g. .wmf, .emf, .tiff). image_alt maps image name
    to its source alt/description text (may be ""). media_count is the
    number of embedded-video (MSO_SHAPE_TYPE.MEDIA) shapes on the slide.
    unresolved_fills is a list of (shape_name, rId) pairs for fill-blips
    whose relationship couldn't be resolved (skipped, not extracted).
    """
    title = None
    title_shape_id = None
    try:
        if slide.shapes.title is not None and slide.shapes.title.text.strip():
            title = slide.shapes.title.text.strip()
            # python-pptx returns a fresh proxy object on every access, so
            # `shape is slide.shapes.title` is always False. Compare the
            # stable shape_id instead so the title shape is actually
            # skipped below, not re-emitted as the slide's first bullet.
            title_shape_id = slide.shapes.title.shape_id
    except (AttributeError, KeyError):
        pass

    body, images, non_web_images, n_img = [], [], [], 0
    image_alt = {}
    media_count = 0
    rid_to_name = {}
    unresolved_fills = []
    for shape in slide.shapes:
        if title_shape_id is not None and shape.shape_id == title_shape_id:
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            n_img += 1
            ext = shape.image.ext
            name = f"slide{idx:02d}-{n_img}.{ext}"
            img_dir.mkdir(parents=True, exist_ok=True)
            (img_dir / name).write_bytes(shape.image.blob)
            images.append(name)
            try:
                descr = shape._element.nvPicPr.cNvPr.get("descr", "") or ""
            except (AttributeError, KeyError):
                descr = ""
            image_alt[name] = " ".join(descr.split())  # collapse embedded newlines
            if ext.lower() not in WEB_SAFE_EXTS:
                non_web_images.append(name)
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
            media_count += 1
            body.append("> \U0001F3AC This slide has an embedded video in "
                         "the original deck (see `original/`).")
        elif getattr(shape, "has_table", False) and shape.has_table:
            body.extend(table_markdown(shape.table))
            body.append("")
        elif shape.has_text_frame:
            body.extend(text_frame_body_lines(shape.text_frame))

        # Fill-embedded images (placeholder picture-fills, autoshape/logo
        # fills, a MEDIA shape's poster frame) -- see _extract_fill_images.
        # Runs for every shape reaching this point, after its own handling
        # (if any) above; a shape with no fill blip is a no-op.
        n_img = _extract_fill_images(shape, slide, idx, img_dir, n_img,
                                      rid_to_name, images, image_alt,
                                      non_web_images, unresolved_fills)

    notes = None
    if slide.has_notes_slide:
        text = slide.notes_slide.notes_text_frame.text.strip()
        notes = text or None
    return (title, body, images, notes, non_web_images, image_alt,
            media_count, unresolved_fills)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("outdir")
    ap.add_argument("--title", required=True)
    ap.add_argument("--week", required=True, type=int)
    ap.add_argument("--topic", required=True)
    args = ap.parse_args()

    src = Path(args.pptx)
    outdir = Path(args.outdir)
    outdir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(src)

    lines = ["---", "marp: true", "theme: default", "paginate: true",
             f'title: "{args.title}"', f"week: {args.week}",
             f"topic: {args.topic}", "type: lecture",
             f'source: "{src.name}"', "---", ""]
    warnings, total_images = [], 0

    for idx, slide in enumerate(prs.slides, start=1):
        (title, body, images, notes, non_web_images, image_alt, media_count,
         unresolved_fills) = extract_slide(slide, idx, outdir / "img")
        total_images += len(images)
        if idx > 1:
            lines += ["---", ""]
        heading = "# " if idx == 1 else "## "
        lines.append(heading + (title or f"Slide {idx}"))
        lines.append("")
        if body:
            lines += body + [""]
        for name in images:
            lines.append(f"![{image_alt.get(name, '')}](img/{name})")
            if name in non_web_images:
                lines.append(f"<!-- image img/{name} not web-renderable -->")
                lines.append(f"> ⚠️ `img/{name}` is not "
                              "web-renderable — convert to PNG manually.")
            lines.append("")
        if notes:
            lines += ["<!-- Speaker notes:", notes, "-->", ""]
        if not title and not body and not images:
            warnings.append(f"slide {idx}: no extractable content")
        for name in non_web_images:
            ext = name.rsplit(".", 1)[-1]
            warnings.append(f"slide {idx}: image saved as .{ext} "
                             "- not web-renderable, needs manual conversion")
        for _ in range(media_count):
            warnings.append(f"slide {idx}: embedded video (MEDIA shape) "
                             "not extracted - see original/ for the source pptx")
        for shape_name, rid in unresolved_fills:
            warnings.append(f"slide {idx}: fill-embedded image on shape "
                             f"{shape_name!r} (rId {rid}) could not be "
                             "resolved - skipped")

    (outdir / "slides.md").write_text("\n".join(lines) + "\n", encoding="utf-8")
    for w in warnings:
        print("WARN:", w, file=sys.stderr)
    print(f"OK: {len(prs.slides)} slides, {total_images} images, "
          f"{len(warnings)} warnings")


if __name__ == "__main__":
    main()
